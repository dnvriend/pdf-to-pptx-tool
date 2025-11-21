"""PDF to PPTX converter module.

This module provides functionality to convert PDF documents to PowerPoint
presentations by converting each PDF page to an image and embedding it in a slide.

Note: This code was generated with assistance from AI coding tools
and has been reviewed and tested by a human.
"""

from pathlib import Path

from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

from pdf_to_pptx_tool.logging_config import get_logger

logger = get_logger(__name__)


def convert_pdf_to_pptx(input_pdf: str, output_pptx: str, dpi: int = 200) -> None:
    """Convert a PDF file to PPTX format.

    Converts each page of the PDF to a high-resolution image and embeds it
    in a PowerPoint slide. Images are sized to fit the slide dimensions (16:9).

    Args:
        input_pdf: Path to input PDF file
        output_pptx: Path to output PPTX file
        dpi: Resolution for PDF page conversion (default: 200)

    Raises:
        FileNotFoundError: If input PDF file does not exist
        PermissionError: If unable to write output file
        RuntimeError: If PDF conversion fails
    """
    input_path = Path(input_pdf)

    # Validate input file exists
    logger.debug("Validating input file: %s", input_pdf)
    if not input_path.exists():
        raise FileNotFoundError(f"Input PDF file not found: {input_pdf}")

    if not input_path.is_file():
        raise ValueError(f"Input path is not a file: {input_pdf}")

    file_size_mb = input_path.stat().st_size / (1024 * 1024)
    logger.debug("Input file size: %.2f MB", file_size_mb)
    logger.info("Converting %s to %s (DPI: %d)", input_pdf, output_pptx, dpi)

    # Convert PDF pages to images
    logger.info("Converting PDF pages to images...")
    logger.debug("Using DPI setting: %d", dpi)
    try:
        images = convert_from_path(str(input_path), dpi=dpi)
        page_count = len(images)
        logger.debug("Successfully converted %d pages to images", page_count)
    except Exception as e:
        logger.error("PDF conversion failed: %s", type(e).__name__)
        logger.debug("Full traceback:", exc_info=True)
        raise RuntimeError(f"Failed to convert PDF pages: {e}") from e

    logger.info("Converted %d pages", page_count)

    # Create PowerPoint presentation
    logger.info("Creating PowerPoint presentation...")
    prs = Presentation()

    # Set slide dimensions (16:9 aspect ratio)
    slide_width = Inches(10)
    slide_height = Inches(5.625)
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    logger.debug('Slide dimensions: 10.00" x 5.62" (16:9 aspect ratio)')

    # Add each image as a slide
    for i, image in enumerate(images, 1):
        logger.debug("Processing slide %d/%d", i, page_count)

        # Save image temporarily
        temp_image_path = f"temp_page_{i}.png"
        logger.debug("Saving temporary image: %s", temp_image_path)
        image.save(temp_image_path, "PNG")

        try:
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            logger.debug("Added blank slide %d", i)

            # Add image to slide (fit to slide dimensions)
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(
                temp_image_path, left, top, width=slide_width, height=slide_height
            )
            logger.debug("Added image to slide %d", i)
        finally:
            # Clean up temporary image
            Path(temp_image_path).unlink(missing_ok=True)
            logger.debug("Cleaned up temporary image: %s", temp_image_path)

    # Save presentation
    logger.info("Saving presentation to %s", output_pptx)
    try:
        prs.save(output_pptx)
        output_path = Path(output_pptx)
        output_size_mb = output_path.stat().st_size / (1024 * 1024)
        logger.debug("Output file size: %.2f MB", output_size_mb)
        logger.info("Successfully converted to %s", output_pptx)
    except Exception as e:
        logger.error("Failed to save presentation: %s", type(e).__name__)
        logger.debug("Full traceback:", exc_info=True)
        raise
